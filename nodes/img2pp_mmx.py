# ~/ComfyUI/custom_nodes/Aiya_mmx/nodes/img2pp_mmx.py
from __future__ import annotations
import os
import uuid
import re
import time
import fcntl
from pathlib import Path

import torch
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

import folder_paths
from ..register import register_node

# --------------------------------------------------
#  通用工具：把任意 Tensor IMAGE → PIL.Image
# --------------------------------------------------
def tensor2pil(tensor):
    if tensor.ndim == 3:
        tensor = tensor.unsqueeze(0)
    tensor = tensor.cpu()
    tensor = torch.clamp(tensor, 0, 1)
    imgs = []
    for im in tensor:
        im = (im.numpy() * 255).astype(np.uint8)
        imgs.append(Image.fromarray(im))
    return imgs

# --------------------------------------------------
#  核心节点：Img2PdfPpt_mmx
# --------------------------------------------------
class Img2PdfPpt_mmx:
    DESCRIPTION = (
        "📄 把多张 IMAGE 一键合并成 PDF + PPTX\n"
        "• 支持 batch 或多路插口\n"
        "• 每节点每次递增，互不干扰\n"
        "• 支持子目录/日期变量"
    )

    def __init__(self):
        self.output_dir = folder_paths.get_output_directory()
        self.type = "output"
        self.prefix_append = ""

    @classmethod
    def INPUT_TYPES(cls):
        return {
            "required": {
                "filename_prefix": ("STRING", {"default": "ComfyUI"}),
                "subfolder": ("STRING", {"default": "", "multiline": False}),
            },
            "optional": {
                f"image_{i}": ("IMAGE",)
                for i in range(1, 10)
            },
            "hidden": {
                "prompt": "PROMPT",
                "extra_pnginfo": "EXTRA_PNGINFO"
            }
        }

    RETURN_TYPES = ("STRING", "STRING")
    RETURN_NAMES = ("pdf_path", "pptx_path")
    FUNCTION = "convert"
    OUTPUT_NODE = True
    CATEGORY = "哎呀✦MMX/图像"

    # ---------- 清理文件名前缀 ----------
    def _clean_filename_prefix(self, prefix: str) -> str:
        """清理文件名前缀，移除非法字符"""
        if not prefix:
            return "ComfyUI"
        
        # 移除首尾空格
        clean = prefix.strip()
        
        # 将多个空格替换为单个下划线
        clean = re.sub(r'\s+', '_', clean)
        
        # 移除Windows/Linux文件系统中不允许的字符
        # 保留字母、数字、中文、下划线、短横线、点
        clean = re.sub(r'[<>:"\\|?*\x00-\x1f]', '', clean)
        
        # 确保不以点或空格开头或结尾
        clean = clean.strip('. ')
        
        # 如果清理后为空，返回默认值
        if not clean:
            clean = "ComfyUI"
            
        return clean

    # ---------- 清理子目录路径 ----------
    def _clean_subfolder(self, subfolder: str) -> str:
        """清理子目录路径，确保安全"""
        if not subfolder:
            return ""
        
        # 移除首尾空格和斜杠
        clean = subfolder.strip().strip('/\\')
        
        # 分割路径部分
        parts = []
        for part in clean.split('/'):
            if part:
                # 清理每个部分
                part_clean = re.sub(r'[<>:"|?*\x00-\x1f]', '', part.strip())
                if part_clean:
                    parts.append(part_clean)
        
        # 重新组合
        if parts:
            return '/'.join(parts)
        return ""

    # ---------- 使用文件锁确保线程安全的计数器 ----------
    def _get_next_counter(self, output_folder: Path, prefix: str):
        """使用文件锁确保线程安全的计数器"""
        lock_file = output_folder / f".{prefix}_counter.lock"
        max_attempts = 30
        attempt = 0
        
        while attempt < max_attempts:
            try:
                # 尝试获取文件锁
                lock_fd = os.open(str(lock_file), os.O_CREAT | os.O_RDWR)
                fcntl.flock(lock_fd, fcntl.LOCK_EX | fcntl.LOCK_NB)
                
                try:
                    # 获取锁后，检查所有现有文件
                    existing_numbers = set()
                    
                    # PDF文件
                    for pdf_file in output_folder.glob(f"{prefix}_*.pdf"):
                        match = re.match(rf"^{re.escape(prefix)}_(\d{{5}})\.pdf$", pdf_file.name)
                        if match:
                            try:
                                existing_numbers.add(int(match.group(1)))
                            except ValueError:
                                continue
                    
                    # PPTX文件
                    for pptx_file in output_folder.glob(f"{prefix}_*.pptx"):
                        match = re.match(rf"^{re.escape(prefix)}_(\d{{5}})\.pptx$", pptx_file.name)
                        if match:
                            try:
                                existing_numbers.add(int(match.group(1)))
                            except ValueError:
                                continue
                    
                    # 计算下一个编号
                    if not existing_numbers:
                        next_num = 1
                    else:
                        max_num = max(existing_numbers)
                        # 检查空缺编号
                        for i in range(1, max_num + 1):
                            if i not in existing_numbers:
                                next_num = i
                                break
                        else:
                            next_num = max_num + 1
                    
                    return next_num
                    
                finally:
                    # 释放锁
                    fcntl.flock(lock_fd, fcntl.LOCK_UN)
                    os.close(lock_fd)
                    # 删除锁文件
                    try:
                        os.unlink(lock_file)
                    except:
                        pass
                    
            except (IOError, BlockingIOError):
                # 获取锁失败，等待后重试
                attempt += 1
                time.sleep(0.05)
                continue
            except Exception as e:
                print(f"获取文件锁时出错: {e}")
                break
        
        # 后备方案
        import random
        fallback_num = int(time.time() * 1000) % 1000000
        print(f"警告: 无法获取文件锁，使用后备编号: {fallback_num}")
        return fallback_num

    def convert(self, filename_prefix="ComfyUI", subfolder="", prompt=None, extra_pnginfo=None, **kwargs):
        # 1. 收集所有非空 IMAGE
        images = []
        for k in kwargs:
            if k.startswith("image_") and kwargs[k] is not None:
                images.append(kwargs[k])
        if not images:
            raise RuntimeError("Img2PdfPpt_mmx: 未收到任何图片输入！")

        # 2. 全部转 PIL
        pil_list = []
        for tensor in images:
            pil_list.extend(tensor2pil(tensor))

        # 3. 清理子目录路径
        clean_subfolder = self._clean_subfolder(subfolder)
        
        # 4. 替换日期变量和清理文件名前缀
        from ..date_variable import replace_date_vars
        
        # 替换日期变量
        raw_prefix = replace_date_vars(filename_prefix)
        
        # 清理文件名前缀
        clean_prefix = self._clean_filename_prefix(raw_prefix)
        name_prefix = clean_prefix + self.prefix_append
        name_prefix = self._clean_filename_prefix(name_prefix)
        
        # 5. 手动构建完整输出路径
        # 首先获取基础输出目录
        base_output_dir = Path(self.output_dir)
        
        # 如果有子目录，添加到路径中
        if clean_subfolder:
            # 替换子目录中的日期变量
            clean_subfolder = replace_date_vars(clean_subfolder)
            # 再次清理
            clean_subfolder = self._clean_subfolder(clean_subfolder)
            output_dir = base_output_dir / clean_subfolder
        else:
            output_dir = base_output_dir
        
        # 确保目录存在
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # 6. 获取下一个计数器编号
        counter = self._get_next_counter(output_dir, name_prefix)
        
        # 7. 构建最终文件名和路径
        pdf_file  = f"{name_prefix}_{counter:05}.pdf"
        pptx_file = f"{name_prefix}_{counter:05}.pptx"
        pdf_path  = output_dir / pdf_file
        pptx_path = output_dir / pptx_file
        
        # 对于UI显示，需要计算相对路径
        if clean_subfolder:
            subfolder_for_ui = clean_subfolder
        else:
            # 获取相对于基础输出目录的路径
            try:
                subfolder_for_ui = str(output_dir.relative_to(base_output_dir))
                if subfolder_for_ui == ".":
                    subfolder_for_ui = ""
            except:
                subfolder_for_ui = ""

        # 8. 调试输出
        print(f"Img2PdfPpt_mmx: 文件名前缀: {name_prefix}")
        print(f"Img2PdfPpt_mmx: 子目录: {clean_subfolder}")
        print(f"Img2PdfPpt_mmx: 输出目录: {output_dir}")
        print(f"Img2PdfPpt_mmx: 计数器: {counter}")
        print(f"Img2PdfPpt_mmx: 生成文件 - PDF: {pdf_file}, PPTX: {pptx_file}")

        # 9. 写 PDF
        try:
            pil_list[0].save(
                pdf_path,
                "PDF",
                quality=95,
                optimize=True,
                append_images=pil_list[1:],
                save_all=True
            )
        except Exception as e:
            print(f"Img2PdfPpt_mmx: 保存PDF时出错: {e}")
            # 确保目录存在
            pdf_path.parent.mkdir(parents=True, exist_ok=True)
            pil_list[0].save(
                pdf_path,
                "PDF",
                quality=95,
                optimize=True,
                append_images=pil_list[1:],
                save_all=True
            )

        # 10. 写 PPTX - 保持原始图片比例
        prs = Presentation()
        
        # PPTX幻灯片的标准尺寸（16:9）
        slide_width = Inches(10)  # 10英寸宽
        slide_height = Inches(5.625)  # 5.625英寸高（16:9比例）
        
        for img in pil_list:
            # 创建新幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白幻灯片
            
            # 获取原始图片尺寸
            img_width_px, img_height_px = img.size
            img_ratio = img_width_px / img_height_px
            
            # 保存临时图片
            tmp = output_dir / f"_tmp_{uuid.uuid4().hex}.png"
            img.save(tmp)
            
            # 根据图片比例调整尺寸
            slide_ratio = slide_width / slide_height
            
            if img_ratio > slide_ratio:
                # 图片比幻灯片宽（横向），以宽度为准
                width = slide_width
                height = width / img_ratio
                # 垂直居中
                top = (slide_height - height) / 2
                left = 0
            else:
                # 图片比幻灯片高（纵向），以高度为准
                height = slide_height
                width = height * img_ratio
                # 水平居中
                left = (slide_width - width) / 2
                top = 0
            
            # 添加图片到幻灯片，保持原始比例
            slide.shapes.add_picture(
                str(tmp),
                left,
                top,
                width=width,
                height=height
            )
            
            # 删除临时文件
            tmp.unlink(missing_ok=True)
        
        try:
            prs.save(pptx_path)
        except Exception as e:
            print(f"Img2PdfPpt_mmx: 保存PPTX时出错: {e}")
            # 确保目录存在
            pptx_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(pptx_path)

        # 11. 返回
        return {"ui": {"images": [{"filename": pdf_file,  "subfolder": subfolder_for_ui, "type": self.type},
                                  {"filename": pptx_file, "subfolder": subfolder_for_ui, "type": self.type}]},
                "result": (str(pdf_path), str(pptx_path))}

register_node(Img2PdfPpt_mmx, "Img2PdfPpt_mmx")